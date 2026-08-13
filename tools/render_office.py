#!/usr/bin/env python3
"""Render compact course JSON into native DOCX and PPTX files."""

from __future__ import annotations

import argparse
import copy
import json
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any, Iterable


FONT_NAME = "Microsoft YaHei"
BLUE = "0083FE"
TEAL = "0BB8B8"
CYAN = "00FFF0"
WHITE = "FFFFFF"
CHARCOAL = "333333"
LIGHT_BLUE = "EAF5FF"
LIGHT_GRAY = "F2F4F7"
MID_GRAY = "A9B0B8"
SUPPORTED_SLIDE_TYPES = {"cover", "bullets", "content", "activity", "question"}


class RenderError(ValueError):
    """Raised when renderer input or the runtime is invalid."""


def load_payload(path: Path) -> dict[str, Any]:
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise RenderError(f"无法读取 JSON：{path}: {exc}") from exc
    if not isinstance(payload, dict):
        raise RenderError("JSON 根节点必须是对象")
    return payload


def _documents(payload: dict[str, Any]) -> list[dict[str, Any]]:
    documents = payload.get("documents", [])
    if "document" in payload:
        documents = [payload["document"], *documents]
    if not isinstance(documents, list) or not all(isinstance(item, dict) for item in documents):
        raise RenderError("documents 必须是对象数组")
    return documents


def _require_text(value: Any, field: str) -> str:
    if not isinstance(value, str) or not value.strip():
        raise RenderError(f"{field} 必须是非空文本")
    return value.strip()


def _safe_filename(value: Any, field: str, suffix: str) -> str:
    filename = _require_text(value, field)
    path = Path(filename)
    if path.name != filename or filename in {".", ".."}:
        raise RenderError(f"{field} 只能是文件名，不能包含目录")
    if path.suffix.lower() != suffix:
        raise RenderError(f"{field} 必须以 {suffix} 结尾")
    return filename


def _compact_length(text: str) -> int:
    return len("".join(text.split()))


def validate_payload(payload: dict[str, Any]) -> None:
    if payload.get("schema_version") != 1:
        raise RenderError("schema_version 必须为 1")
    documents = _documents(payload)
    presentation = payload.get("presentation")
    if not documents and presentation is None:
        raise RenderError("至少提供 documents 或 presentation")

    for doc_index, document in enumerate(documents, start=1):
        prefix = f"documents[{doc_index}]"
        _safe_filename(document.get("filename"), f"{prefix}.filename", ".docx")
        _require_text(document.get("title"), f"{prefix}.title")
        if document.get("audience", "teacher") not in {"teacher", "student"}:
            raise RenderError(f"{prefix}.audience 只能是 teacher 或 student")
        sections = document.get("sections", [])
        if not isinstance(sections, list):
            raise RenderError(f"{prefix}.sections 必须是数组")
        for section_index, section in enumerate(sections, start=1):
            if not isinstance(section, dict):
                raise RenderError(f"{prefix}.sections[{section_index}] 必须是对象")
            questions = section.get("questions", [])
            if not isinstance(questions, list):
                raise RenderError(f"{prefix}.sections[{section_index}].questions 必须是数组")
            for question_index, question in enumerate(questions, start=1):
                _validate_question(
                    question,
                    f"{prefix}.sections[{section_index}].questions[{question_index}]",
                )

    if presentation is not None:
        if not isinstance(presentation, dict):
            raise RenderError("presentation 必须是对象")
        _safe_filename(presentation.get("filename"), "presentation.filename", ".pptx")
        slides = presentation.get("slides", [])
        if not isinstance(slides, list) or not slides:
            raise RenderError("presentation.slides 必须是非空数组")
        for index, slide in enumerate(slides, start=1):
            if not isinstance(slide, dict):
                raise RenderError(f"presentation.slides[{index}] 必须是对象")
            slide_type = slide.get("type", "content")
            if slide_type not in SUPPORTED_SLIDE_TYPES:
                raise RenderError(f"presentation.slides[{index}].type 不受支持：{slide_type}")
            _require_text(slide.get("title"), f"presentation.slides[{index}].title")
            bullets = slide.get("bullets", [])
            if not isinstance(bullets, list) or not all(isinstance(item, str) for item in bullets):
                raise RenderError(f"presentation.slides[{index}].bullets 必须是文本数组")
            if len(bullets) > 3:
                raise RenderError(f"presentation.slides[{index}] 超过 3 个要点，请拆成多个逻辑页")
            for bullet_index, bullet in enumerate(bullets, start=1):
                if _compact_length(bullet) > 36:
                    raise RenderError(
                        f"presentation.slides[{index}].bullets[{bullet_index}] 超过 36 个非空白字符"
                    )
            if slide.get("build", "cumulative") not in {"cumulative", "none"}:
                raise RenderError(f"presentation.slides[{index}].build 只能是 cumulative 或 none")


def _validate_question(question: Any, field: str) -> None:
    if not isinstance(question, dict):
        raise RenderError(f"{field} 必须是对象")
    _require_text(question.get("prompt"), f"{field}.prompt")
    options = question.get("options")
    if not isinstance(options, list) or len(options) != 4:
        raise RenderError(f"{field}.options 必须正好包含 4 个选项")
    labels: list[str] = []
    scores: list[int] = []
    for option_index, option in enumerate(options, start=1):
        if not isinstance(option, dict):
            raise RenderError(f"{field}.options[{option_index}] 必须是对象")
        labels.append(_require_text(option.get("label"), f"{field}.options[{option_index}].label"))
        _require_text(option.get("text"), f"{field}.options[{option_index}].text")
        if "score" in option:
            score = option["score"]
            if not isinstance(score, int) or not 1 <= score <= 4:
                raise RenderError(f"{field}.options[{option_index}].score 必须是 1-4 的整数")
            scores.append(score)
    if labels != ["A", "B", "C", "D"]:
        raise RenderError(f"{field} 的选项标签必须依次为 A、B、C、D")
    if scores and (len(scores) != 4 or scores.count(4) != 1):
        raise RenderError(f"{field} 必须为每个选项赋分，且只能有一个 4 分选项")


def expand_slides(slides: list[dict[str, Any]], build_mode: str) -> list[dict[str, Any]]:
    expanded: list[dict[str, Any]] = []
    for logical_index, source in enumerate(slides, start=1):
        slide = copy.deepcopy(source)
        bullets = slide.get("bullets", [])
        cumulative = (
            build_mode == "live"
            and slide.get("build", "cumulative") == "cumulative"
            and len(bullets) > 1
        )
        total = len(bullets) if cumulative else 1
        for step in range(1, total + 1):
            state = copy.deepcopy(slide)
            if cumulative:
                state["bullets"] = bullets[:step]
                state["_current_bullet"] = step - 1
            state["_logical_index"] = logical_index
            state["_build_step"] = step
            state["_build_total"] = total
            expanded.append(state)
    return expanded


def _office_imports() -> dict[str, Any]:
    try:
        from docx import Document
        from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import Cm, Pt, RGBColor as DocxRGBColor
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor as PptxRGBColor
        from pptx.util import Inches, Pt as PptxPt
    except ImportError as exc:
        raise RenderError(
            "缺少 Office 依赖，请先运行：python -m pip install -r tools/requirements.txt"
        ) from exc
    return locals()


def _set_docx_run_font(run: Any, imports: dict[str, Any], size: float | None = None) -> None:
    run.font.name = FONT_NAME
    if size is not None:
        run.font.size = imports["Pt"](size)
    r_pr = run._element.get_or_add_rPr()
    r_fonts = r_pr.rFonts
    if r_fonts is None:
        r_fonts = imports["OxmlElement"]("w:rFonts")
        r_pr.insert(0, r_fonts)
    r_fonts.set(imports["qn"]("w:eastAsia"), FONT_NAME)


def _set_docx_style_font(style: Any, imports: dict[str, Any], size: float, color: str) -> None:
    style.font.name = FONT_NAME
    style.font.size = imports["Pt"](size)
    style.font.color.rgb = imports["DocxRGBColor"].from_string(color)
    r_pr = style._element.get_or_add_rPr()
    r_fonts = r_pr.rFonts
    if r_fonts is None:
        r_fonts = imports["OxmlElement"]("w:rFonts")
        r_pr.insert(0, r_fonts)
    r_fonts.set(imports["qn"]("w:eastAsia"), FONT_NAME)


def _add_hyperlink(paragraph: Any, text: str, url: str, imports: dict[str, Any]) -> None:
    part = paragraph.part
    relationship_id = part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = imports["OxmlElement"]("w:hyperlink")
    hyperlink.set(imports["qn"]("r:id"), relationship_id)
    run = imports["OxmlElement"]("w:r")
    run_properties = imports["OxmlElement"]("w:rPr")
    color = imports["OxmlElement"]("w:color")
    color.set(imports["qn"]("w:val"), BLUE)
    underline = imports["OxmlElement"]("w:u")
    underline.set(imports["qn"]("w:val"), "single")
    fonts = imports["OxmlElement"]("w:rFonts")
    fonts.set(imports["qn"]("w:ascii"), FONT_NAME)
    fonts.set(imports["qn"]("w:eastAsia"), FONT_NAME)
    run_properties.extend([fonts, color, underline])
    run.append(run_properties)
    text_element = imports["OxmlElement"]("w:t")
    text_element.text = text
    run.append(text_element)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _set_repeat_header(row: Any, imports: dict[str, Any]) -> None:
    tr_pr = row._tr.get_or_add_trPr()
    header = imports["OxmlElement"]("w:tblHeader")
    header.set(imports["qn"]("w:val"), "true")
    tr_pr.append(header)


def _prevent_row_split(row: Any, imports: dict[str, Any]) -> None:
    tr_pr = row._tr.get_or_add_trPr()
    tr_pr.append(imports["OxmlElement"]("w:cantSplit"))


def _shade_cell(cell: Any, color: str, imports: dict[str, Any]) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shading = imports["OxmlElement"]("w:shd")
    shading.set(imports["qn"]("w:fill"), color)
    tc_pr.append(shading)


def _add_page_number(paragraph: Any, imports: dict[str, Any]) -> None:
    paragraph.alignment = imports["WD_ALIGN_PARAGRAPH"].CENTER
    run = paragraph.add_run("第 ")
    _set_docx_run_font(run, imports, 9)
    begin = imports["OxmlElement"]("w:fldChar")
    begin.set(imports["qn"]("w:fldCharType"), "begin")
    instruction = imports["OxmlElement"]("w:instrText")
    instruction.set(imports["qn"]("xml:space"), "preserve")
    instruction.text = "PAGE"
    end = imports["OxmlElement"]("w:fldChar")
    end.set(imports["qn"]("w:fldCharType"), "end")
    run._r.extend([begin, instruction, end])
    tail = paragraph.add_run(" 页")
    _set_docx_run_font(tail, imports, 9)


def _add_document_table(document: Any, table_data: dict[str, Any], imports: dict[str, Any]) -> None:
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if not headers or not isinstance(rows, list):
        raise RenderError("文档表格必须提供 headers 和 rows")
    if any(not isinstance(row, list) or len(row) != len(headers) for row in rows):
        raise RenderError("文档表格每行列数必须与 headers 一致")
    table = document.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    table.alignment = imports["WD_TABLE_ALIGNMENT"].CENTER
    for index, header_text in enumerate(headers):
        cell = table.rows[0].cells[index]
        cell.text = str(header_text)
        _shade_cell(cell, BLUE, imports)
        cell.vertical_alignment = imports["WD_CELL_VERTICAL_ALIGNMENT"].CENTER
        for run in cell.paragraphs[0].runs:
            _set_docx_run_font(run, imports, 9)
            run.font.bold = True
            run.font.color.rgb = imports["DocxRGBColor"].from_string(WHITE)
    _set_repeat_header(table.rows[0], imports)
    for row_values in rows:
        row = table.add_row()
        _prevent_row_split(row, imports)
        for index, value in enumerate(row_values):
            cell = row.cells[index]
            cell.text = str(value)
            cell.vertical_alignment = imports["WD_CELL_VERTICAL_ALIGNMENT"].CENTER
            for run in cell.paragraphs[0].runs:
                _set_docx_run_font(run, imports, 9)


def _add_sources(document: Any, sources: list[Any], imports: dict[str, Any]) -> None:
    for index, source in enumerate(sources, start=1):
        if not isinstance(source, dict):
            raise RenderError("sources 中每一项必须是对象")
        title = _require_text(source.get("title"), f"sources[{index}].title")
        publisher = str(source.get("publisher", "发布方未标明"))
        paragraph = document.add_paragraph()
        heading = paragraph.add_run(f"{index}. {title}｜{publisher}")
        heading.bold = True
        _set_docx_run_font(heading, imports, 10.5)
        url = source.get("url")
        if isinstance(url, str) and url.strip():
            link_paragraph = document.add_paragraph()
            _add_hyperlink(link_paragraph, url.strip(), url.strip(), imports)
        for label, key in (
            ("日期", "date"),
            ("访问日期", "accessed"),
            ("关键信息", "key_points"),
            ("本课采用", "application"),
            ("限制/待核实", "caveat"),
            ("许可/使用条件", "license"),
        ):
            value = source.get(key)
            if value:
                item = document.add_paragraph()
                label_run = item.add_run(f"{label}：")
                label_run.bold = True
                _set_docx_run_font(label_run, imports, 10)
                value_run = item.add_run(str(value))
                _set_docx_run_font(value_run, imports, 10)


def _add_questions(
    document: Any,
    questions: list[dict[str, Any]],
    audience: str,
    imports: dict[str, Any],
) -> None:
    for index, question in enumerate(questions, start=1):
        prompt = document.add_paragraph()
        prompt.paragraph_format.space_before = imports["Pt"](8)
        prompt_run = prompt.add_run(f"{index}. {question['prompt']}")
        prompt_run.bold = True
        _set_docx_run_font(prompt_run, imports, 11)
        for option in question["options"]:
            text = f"{option['label']}. {option['text']}"
            if audience == "teacher" and "score" in option:
                level = option.get("level", "")
                rationale = option.get("rationale", "")
                text += f"（{option['score']} 分{('，' + level) if level else ''}）"
                if rationale:
                    text += f" {rationale}"
            paragraph = document.add_paragraph(style="List Bullet")
            run = paragraph.add_run(text)
            _set_docx_run_font(run, imports, 10.5)
            if audience == "teacher" and option.get("score") == 4:
                run.bold = True
                run.font.color.rgb = imports["DocxRGBColor"].from_string(BLUE)


def render_docx(document_data: dict[str, Any], output_dir: Path, imports: dict[str, Any]) -> Path:
    document = imports["Document"]()
    section = document.sections[0]
    section.page_width = imports["Cm"](21)
    section.page_height = imports["Cm"](29.7)
    section.top_margin = imports["Cm"](1.6)
    section.bottom_margin = imports["Cm"](1.6)
    section.left_margin = imports["Cm"](1.7)
    section.right_margin = imports["Cm"](1.7)

    styles = document.styles
    normal = styles["Normal"]
    _set_docx_style_font(normal, imports, 10.5, CHARCOAL)
    for style_name, size, color in (
        ("Title", 22, BLUE),
        ("Heading 1", 16, BLUE),
        ("Heading 2", 13, TEAL),
    ):
        _set_docx_style_font(styles[style_name], imports, size, color)

    title = document.add_heading(document_data["title"], 0)
    title.alignment = imports["WD_ALIGN_PARAGRAPH"].CENTER
    subtitle = document_data.get("subtitle")
    if subtitle:
        paragraph = document.add_paragraph()
        paragraph.alignment = imports["WD_ALIGN_PARAGRAPH"].CENTER
        run = paragraph.add_run(str(subtitle))
        run.font.color.rgb = imports["DocxRGBColor"].from_string(CHARCOAL)
        _set_docx_run_font(run, imports, 11)

    audience = document_data.get("audience", "teacher")
    for section_index, section_data in enumerate(document_data.get("sections", [])):
        if section_data.get("page_break_before") and (section_index > 0 or len(document.paragraphs) > 1):
            document.add_page_break()
        if section_data.get("title"):
            document.add_heading(str(section_data["title"]), level=1)
        for text in section_data.get("paragraphs", []):
            paragraph = document.add_paragraph()
            run = paragraph.add_run(str(text))
            _set_docx_run_font(run, imports, 10.5)
        for text in section_data.get("bullets", []):
            paragraph = document.add_paragraph(style="List Bullet")
            run = paragraph.add_run(str(text))
            _set_docx_run_font(run, imports, 10.5)
        if section_data.get("table"):
            _add_document_table(document, section_data["table"], imports)
        if section_data.get("sources"):
            _add_sources(document, section_data["sources"], imports)
        if section_data.get("questions"):
            _add_questions(document, section_data["questions"], audience, imports)

    _add_page_number(section.footer.paragraphs[0], imports)
    document.core_properties.title = document_data["title"]
    document.core_properties.author = "IBL Course Designer"
    output_path = output_dir / document_data["filename"]
    document.save(output_path)
    imports["Document"](output_path)
    return output_path


def _pptx_color(imports: dict[str, Any], value: str) -> Any:
    return imports["PptxRGBColor"].from_string(value)


def _set_pptx_text(
    shape: Any,
    text: str,
    imports: dict[str, Any],
    *,
    size: float,
    color: str = CHARCOAL,
    bold: bool = False,
    alignment: Any | None = None,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = imports["MSO_ANCHOR"].MIDDLE
    paragraph = frame.paragraphs[0]
    if alignment is not None:
        paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = imports["PptxPt"](size)
    run.font.bold = bold
    run.font.color.rgb = _pptx_color(imports, color)


def _add_slide_title(slide: Any, title: str, imports: dict[str, Any]) -> None:
    shape = slide.shapes.add_textbox(
        imports["Inches"](0.75),
        imports["Inches"](0.35),
        imports["Inches"](11.85),
        imports["Inches"](0.75),
    )
    _set_pptx_text(shape, title, imports, size=30, color=BLUE, bold=True)
    line = slide.shapes.add_shape(
        imports["MSO_SHAPE"].RECTANGLE,
        imports["Inches"](0.75),
        imports["Inches"](1.12),
        imports["Inches"](1.2),
        imports["Inches"](0.06),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _pptx_color(imports, TEAL)
    line.line.fill.background()


def _add_footer(slide: Any, footer: str, state: dict[str, Any], imports: dict[str, Any]) -> None:
    build = ""
    if state.get("_build_total", 1) > 1:
        build = f" · {state['_build_step']}/{state['_build_total']}"
    text = f"{footer} · 逻辑页 {state['_logical_index']}{build}"
    shape = slide.shapes.add_textbox(
        imports["Inches"](0.75),
        imports["Inches"](7.05),
        imports["Inches"](11.85),
        imports["Inches"](0.25),
    )
    _set_pptx_text(shape, text, imports, size=10, color=MID_GRAY)


def _add_bullets_slide(slide: Any, state: dict[str, Any], imports: dict[str, Any]) -> None:
    bullets = state.get("bullets", [])
    current = state.get("_current_bullet", len(bullets) - 1)
    top = 1.55
    spacing = 1.55 if len(bullets) <= 3 else 1.25
    for index, bullet in enumerate(bullets):
        is_current = index == current
        number = slide.shapes.add_shape(
            imports["MSO_SHAPE"].OVAL,
            imports["Inches"](1.0),
            imports["Inches"](top + index * spacing),
            imports["Inches"](0.58),
            imports["Inches"](0.58),
        )
        number.fill.solid()
        number.fill.fore_color.rgb = _pptx_color(imports, BLUE if is_current else LIGHT_BLUE)
        number.line.color.rgb = _pptx_color(imports, BLUE)
        _set_pptx_text(
            number,
            str(index + 1),
            imports,
            size=20,
            color=WHITE if is_current else BLUE,
            bold=True,
            alignment=imports["PP_ALIGN"].CENTER,
        )
        text = slide.shapes.add_textbox(
            imports["Inches"](1.9),
            imports["Inches"](top - 0.03 + index * spacing),
            imports["Inches"](9.7),
            imports["Inches"](0.72),
        )
        _set_pptx_text(text, bullet, imports, size=26, bold=is_current)
        divider = slide.shapes.add_shape(
            imports["MSO_SHAPE"].RECTANGLE,
            imports["Inches"](1.9),
            imports["Inches"](top + 0.78 + index * spacing),
            imports["Inches"](9.5),
            imports["Inches"](0.018),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = _pptx_color(imports, LIGHT_GRAY)
        divider.line.fill.background()


def _add_content_slide(slide: Any, state: dict[str, Any], imports: dict[str, Any]) -> None:
    body = str(state.get("body", ""))
    shape = slide.shapes.add_textbox(
        imports["Inches"](1.0),
        imports["Inches"](1.65),
        imports["Inches"](11.25),
        imports["Inches"](4.7),
    )
    _set_pptx_text(shape, body, imports, size=28)


def _add_activity_slide(slide: Any, state: dict[str, Any], imports: dict[str, Any]) -> None:
    fields = [
        ("任务", state.get("task", "")),
        ("时间", state.get("minutes", "")),
        ("方式", state.get("mode", "")),
        ("产出", state.get("output", "")),
    ]
    for index, (label, value) in enumerate(fields):
        y = 1.55 + index * 1.2
        label_shape = slide.shapes.add_textbox(
            imports["Inches"](1.0), imports["Inches"](y), imports["Inches"](1.2), imports["Inches"](0.65)
        )
        _set_pptx_text(label_shape, label, imports, size=19, color=TEAL, bold=True)
        value_shape = slide.shapes.add_textbox(
            imports["Inches"](2.25), imports["Inches"](y), imports["Inches"](9.6), imports["Inches"](0.75)
        )
        _set_pptx_text(value_shape, str(value), imports, size=25, bold=index == 0)


def _add_question_slide(slide: Any, state: dict[str, Any], imports: dict[str, Any]) -> None:
    prompt = slide.shapes.add_textbox(
        imports["Inches"](0.95), imports["Inches"](1.4), imports["Inches"](11.4), imports["Inches"](1.0)
    )
    _set_pptx_text(prompt, str(state.get("prompt", "")), imports, size=25, bold=True)
    options = state.get("options", [])
    for index, option in enumerate(options[:4]):
        column = index % 2
        row = index // 2
        shape = slide.shapes.add_textbox(
            imports["Inches"](1.0 + column * 5.85),
            imports["Inches"](2.65 + row * 1.55),
            imports["Inches"](5.3),
            imports["Inches"](1.05),
        )
        if isinstance(option, dict):
            text = f"{option.get('label', chr(65 + index))}. {option.get('text', '')}"
        else:
            text = f"{chr(65 + index)}. {option}"
        _set_pptx_text(shape, text, imports, size=20)


def render_pptx(
    presentation_data: dict[str, Any],
    output_dir: Path,
    build_mode: str,
    imports: dict[str, Any],
) -> Path:
    presentation = imports["Presentation"]()
    presentation.slide_width = imports["Inches"](13.333)
    presentation.slide_height = imports["Inches"](7.5)
    blank_layout = presentation.slide_layouts[6]
    states = expand_slides(presentation_data["slides"], build_mode)
    footer = str(presentation_data.get("footer", presentation_data.get("title", "课堂课件")))

    for state in states:
        slide = presentation.slides.add_slide(blank_layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = _pptx_color(imports, WHITE)
        slide_type = state.get("type", "content")
        if slide_type == "cover":
            title_shape = slide.shapes.add_textbox(
                imports["Inches"](1.0), imports["Inches"](1.75), imports["Inches"](11.3), imports["Inches"](1.5)
            )
            _set_pptx_text(
                title_shape,
                state["title"],
                imports,
                size=34,
                color=BLUE,
                bold=True,
                alignment=imports["PP_ALIGN"].CENTER,
            )
            subtitle_shape = slide.shapes.add_textbox(
                imports["Inches"](1.5), imports["Inches"](3.45), imports["Inches"](10.3), imports["Inches"](0.85)
            )
            _set_pptx_text(
                subtitle_shape,
                str(state.get("subtitle", presentation_data.get("subtitle", ""))),
                imports,
                size=22,
                color=CHARCOAL,
                alignment=imports["PP_ALIGN"].CENTER,
            )
        else:
            _add_slide_title(slide, state["title"], imports)
            if slide_type == "bullets":
                _add_bullets_slide(slide, state, imports)
            elif slide_type == "activity":
                _add_activity_slide(slide, state, imports)
            elif slide_type == "question":
                _add_question_slide(slide, state, imports)
            else:
                _add_content_slide(slide, state, imports)
        _add_footer(slide, footer, state, imports)
        notes = str(state.get("notes", ""))
        if state.get("_build_total", 1) > 1:
            notes = f"呈现步骤 {state['_build_step']}/{state['_build_total']}。{notes}"
        slide.notes_slide.notes_text_frame.text = notes

    presentation.core_properties.title = str(presentation_data.get("title", "课堂课件"))
    presentation.core_properties.author = "IBL Course Designer"
    output_path = output_dir / presentation_data["filename"]
    presentation.save(output_path)
    reopened = imports["Presentation"](output_path)
    if len(reopened.slides) != len(states):
        raise RenderError("PPTX 保存后的页数检查失败")
    return output_path


def convert_docx_to_pdf(paths: Iterable[Path], output_dir: Path) -> list[Path]:
    executable = shutil.which("soffice") or shutil.which("libreoffice")
    if executable is None:
        raise RenderError("--pdf 需要环境安装 LibreOffice/soffice")
    converted: list[Path] = []
    for path in paths:
        result = subprocess.run(
            [
                executable,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(path),
            ],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
        pdf_path = output_dir / f"{path.stem}.pdf"
        if result.returncode != 0 or not pdf_path.is_file():
            detail = (result.stderr or result.stdout).strip()
            raise RenderError(f"PDF 转换失败：{path.name}: {detail}")
        converted.append(pdf_path)
    return converted


def render(
    payload: dict[str, Any],
    output_dir: Path,
    *,
    only: str = "all",
    build_mode: str = "live",
    pdf: bool = False,
) -> list[Path]:
    validate_payload(payload)
    output_dir.mkdir(parents=True, exist_ok=True)
    imports = _office_imports()
    outputs: list[Path] = []
    docx_paths: list[Path] = []
    if only in {"all", "docx"}:
        for document_data in _documents(payload):
            path = render_docx(document_data, output_dir, imports)
            outputs.append(path)
            docx_paths.append(path)
    if only in {"all", "pptx"} and payload.get("presentation"):
        outputs.append(render_pptx(payload["presentation"], output_dir, build_mode, imports))
    if pdf and docx_paths:
        outputs.extend(convert_docx_to_pdf(docx_paths, output_dir))
    return outputs


def main() -> int:
    for stream in (sys.stdout, sys.stderr):
        reconfigure = getattr(stream, "reconfigure", None)
        if reconfigure is not None:
            reconfigure(encoding="utf-8", errors="replace")
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", type=Path, help="UTF-8 JSON input")
    parser.add_argument("--output-dir", type=Path, default=Path("output"))
    parser.add_argument("--only", choices=("all", "docx", "pptx"), default="all")
    parser.add_argument("--build-mode", choices=("live", "final"), default="live")
    parser.add_argument("--pdf", action="store_true", help="Convert generated DOCX files with LibreOffice")
    args = parser.parse_args()
    try:
        payload = load_payload(args.input)
        outputs = render(
            payload,
            args.output_dir,
            only=args.only,
            build_mode=args.build_mode,
            pdf=args.pdf,
        )
    except (OSError, RenderError, subprocess.SubprocessError) as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    for output in outputs:
        print(f"OK: {output} ({output.stat().st_size} bytes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
